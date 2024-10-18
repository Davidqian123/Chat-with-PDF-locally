import json
import random
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.util import Inches
from collections import defaultdict
from copy import deepcopy
import os
from PIL import Image


class PresentationGenerator:
    def __init__(self, template_path):
        self.template = Presentation(template_path)
        self.slide_types = self._get_slide_types()
        self.output_prs = None

    def _get_slide_types(self):
        slide_types = defaultdict(list)
        for index, slide in enumerate(self.template.slides):
            if slide.has_notes_slide:
                note = slide.notes_slide.notes_text_frame.text.strip()
                note_base = "".join([i for i in note if not i.isdigit()]).lower()
                if note_base in ["first", "chart", "text", "last"]:
                    slide_types[note_base].append({"index": index, "note": note})
        return dict(slide_types)

    def read_json_data(self, file_path):
        with open(file_path, "r") as file:
            return json.load(file)

    def _choose_slide_index(self, slide_data):
        if "ppt_title" in slide_data:
            return random.choice(
                self.slide_types.get("first", [{"index": 0, "note": "first"}])
            )
        elif slide_data.get("image_path"):
            return random.choice(
                self.slide_types.get("chart", [{"index": 1, "note": "chart1"}])
            )
        else:
            return random.choice(
                self.slide_types.get("text", [{"index": 4, "note": "text1"}])
            )

    def _fill_placeholder(self, shape, content):
        if shape.has_text_frame:
            text_frame = shape.text_frame

            # Store original properties
            original_paragraph = text_frame.paragraphs[0]
            original_font = original_paragraph.runs[0].font
            original_alignment = original_paragraph.alignment

            # Clear existing text (this removes all runs)
            text_frame.clear()

            # Split content into bullet points
            bullet_points = content.split("[b]")

            for i, point in enumerate(bullet_points):
                point = point.strip()
                if point:
                    # Add new paragraph for each point
                    if i == 0:
                        paragraph = text_frame.paragraphs[0]
                    else:
                        paragraph = text_frame.add_paragraph()

                    paragraph.text = point

                    # Apply original alignment to the paragraph
                    paragraph.alignment = original_alignment

                    # Set bullet style for all except the first paragraph (if it doesn't start with [b])
                    if i > 0 or content.strip().startswith("[b]"):
                        paragraph.level = 0
                        paragraph.bullet = True

                    # Apply original font properties to the new text
                    for run in paragraph.runs:
                        run.font.name = original_font.name
                        run.font.size = original_font.size
                        run.font.bold = original_font.bold
                        run.font.italic = original_font.italic

                        # Handle different color types
                        if original_font.color.type == MSO_COLOR_TYPE.RGB:
                            run.font.color.rgb = original_font.color.rgb
                        elif original_font.color.type == MSO_COLOR_TYPE.SCHEME:
                            run.font.color.theme_color = original_font.color.theme_color

    def _create_picture_placeholder(self, slide, shape, image_path):
        # Get the position and dimensions of the original shape
        original_left = shape.left
        original_top = shape.top
        original_width = shape.width
        original_height = shape.height

        # Get the slide height from the presentation
        slide_height = self.output_prs.slide_height

        # Open the image and calculate the aspect ratio
        with Image.open(image_path) as img:
            aspect_ratio = img.height / img.width

        # Calculate new dimensions based on aspect ratio and max width
        new_width = original_width
        new_height = int(new_width * aspect_ratio)

        # Ensure the height doesn't exceed 65% of the slide height
        max_allowed_height = int(slide_height * 0.65)
        if new_height > max_allowed_height:
            new_height = max_allowed_height
            # Recalculate width to maintain aspect ratio
            new_width = int(new_height / aspect_ratio)

        # Calculate horizontal centering offset
        left_offset = (original_width - new_width) / 2

        # Calculate new position (centered horizontally, aligned to top)
        new_left = original_left + left_offset
        new_top = original_top  # Align to the top of the original shape

        # Delete the original shape
        sp = shape._element
        sp.getparent().remove(sp)

        # Add a new picture to the slide, centered horizontally and aligned to the top
        pic = slide.shapes.add_picture(
            image_path, new_left, new_top, new_width, new_height
        )
        return pic

    def _add_content_to_slide(self, slide, slide_data):
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "title" in shape.text.lower():
                    if "ppt_title" in slide_data:
                        self._fill_placeholder(shape, slide_data["ppt_title"])
                    else:
                        self._fill_placeholder(shape, slide_data.get("title_text", ""))
                elif "content" in shape.text.lower():
                    if "sub_title" in slide_data:
                        self._fill_placeholder(shape, slide_data["sub_title"])
                    else:
                        self._fill_placeholder(shape, slide_data.get("main_text", ""))
                elif "image" in shape.text.lower():
                    image_path = slide_data.get("image_path")
                    if image_path:
                        self._create_picture_placeholder(slide, shape, image_path)

    def _copy_slide(self, slide, note):
        new_slide = self.output_prs.slides.add_slide(slide.slide_layout)

        # Add background image based on note
        width = Inches(self.output_prs.slide_width.inches)
        height = Inches(self.output_prs.slide_height.inches)
        background_image_path = f"ppt_background_images/{note}.png"
        if os.path.exists(background_image_path):
            new_slide.shapes.add_picture(
                background_image_path, Inches(0), Inches(0), width, height
            )

        # Copy all shapes from the original slide
        for shape in slide.shapes:
            el = shape.element
            newel = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel, "p:extLst")

        return new_slide

    def generate_presentation(self, json_file_path, output_file_path):
        slides_data = self.read_json_data(json_file_path)

        # Add placeholder for title and subtitle
        ppt_title = "2024 ASUS ZENBOOK S 16"
        sub_title = "Provided by Nexa AI"

        # Insert the title and subtitle into slides_data if not present
        if slides_data and "ppt_title" not in slides_data[0]:
            slides_data.insert(
                0, {"ppt_title": ppt_title, "sub_title": sub_title}
            )
        else:
            slides_data[0]["ppt_title"] = ppt_title
            slides_data[0]["sub_title"] = sub_title

        # Initialize the presentation
        self.output_prs = Presentation()
        self.output_prs.slide_width = self.template.slide_width
        self.output_prs.slide_height = self.template.slide_height

        # Process each slide
        for slide_data in slides_data:
            template_slide_info = self._choose_slide_index(slide_data)
            template_slide_index = template_slide_info["index"]
            template_slide = self.template.slides[template_slide_index]
            new_slide = self._copy_slide(template_slide, template_slide_info["note"])
            self._add_content_to_slide(new_slide, slide_data)

        # Add the last slide if it exists in the template
        if "last" in self.slide_types:
            last_slide_info = random.choice(self.slide_types["last"])
            last_slide_index = last_slide_info["index"]
            last_slide = self.template.slides[last_slide_index]
            self._copy_slide(last_slide, last_slide_info["note"])

        # Save the presentation
        self.output_prs.save(output_file_path)
        print(f"Presentation saved as {output_file_path}")


# Usage
if __name__ == "__main__":
    generator = PresentationGenerator(
        "files/ppt_template.pptx"
    )
    generator.generate_presentation(
        "slides_data.json", "example_presentation_with_charts_and_text.pptx"
    )
    # slide_types = generator._get_slide_types()
    # print(slide_types)
